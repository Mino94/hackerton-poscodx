"""graphics_spec → python-pptx 도형/이미지 — 우선순위: 이미지 > elements > visual_builder fallback."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from autopm.ppt import layout_engine, theme
from autopm.ppt.slide_schema import SlideSpec
from autopm.ppt.visual_builder import apply_visual

_LEFT = Inches(0.55)
_W = Inches(12.67)
# layout_engine.CONTENT_TOP 과 맞춰 제목·콜아웃 아래에 장표를 둔다.
_IMG_TOP = layout_engine.CONTENT_TOP
_IMG_H = Inches(4.38)


def _project_root() -> Path:
    # src/autopm/ppt/graphics_layout.py → parents[3] = repo root
    return Path(__file__).resolve().parents[3]


def add_graphics_elements_horizontal(slide: Slide, elements: list[dict[str, Any]]) -> None:
    """graphics_spec.elements — 수평 rounded_box + arrow 체인."""
    boxes = [e for e in elements if e.get("type") == "rounded_box"]
    if not boxes:
        return
    labels = [str(b.get("label", ""))[:120] for b in boxes]
    layout_engine.add_process_flow(slide, labels)


def add_graphics_elements_vertical(slide: Slide, elements: list[dict[str, Any]]) -> None:
    """arrow_down이 섞인 아키텍처 스택 — 세로로 박스를 쌓는다."""
    y = float(layout_engine.CONTENT_TOP)
    box_w = _W
    box_h = Inches(0.75)
    gap = Inches(0.15)
    x = _LEFT
    for el in elements:
        t = el.get("type")
        if t == "rounded_box":
            shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, int(y), int(box_w), int(box_h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = theme.CARD_FILL
            shp.line.color.rgb = theme.ACCENT
            tf = shp.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = str(el.get("label", ""))[:200]
            p.font.size = Pt(11)
            p.font.color.rgb = theme.BODY_RGB
            p.alignment = PP_ALIGN.CENTER
            y += float(box_h) + float(gap)
        elif t == "arrow_down":
            arr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                x + float(box_w) / 2 - Inches(0.12),
                int(y),
                Inches(0.24),
                Inches(0.2),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = theme.ACCENT
            y += float(Inches(0.28))


def add_graphics_elements_funnel(slide: Slide, elements: list[dict[str, Any]]) -> None:
    """퍼널 — 폭이 줄어드는 사다리꼴을 세로로 배치."""
    traps = [e for e in elements if e.get("type") == "trapezoid"]
    if not traps:
        return
    n = len(traps)
    y = float(layout_engine.CONTENT_TOP)
    max_w = float(_W)
    for i, el in enumerate(traps):
        frac = (n - i) / (n + 1)
        w = max_w * max(0.35, frac)
        off = (max_w - w) / 2
        h = Inches(0.72)
        b = 220 - (i * 12)
        fill = RGBColor(0xDD, max(0xD0, min(0xFF, b)), 0xF5)
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.TRAPEZOID,
            int(_LEFT + off),
            int(y),
            int(w),
            int(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(el.get("label", ""))[:80]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        y += float(h) + float(Inches(0.08))


def add_graphics_elements(slide: Slide, elements: list[dict[str, Any]]) -> None:
    """elements 타입에 따라 적절한 배치 함수로 위임한다."""
    if not elements:
        return
    if any(e.get("type") == "trapezoid" for e in elements):
        add_graphics_elements_funnel(slide, elements)
        return
    if any(e.get("type") == "arrow_down" for e in elements):
        add_graphics_elements_vertical(slide, elements)
        return
    add_graphics_elements_horizontal(slide, elements)


def apply_graphics_or_visual(slide: Slide, slide_spec: SlideSpec) -> None:
    """
    PPT Composer 최종 단계 — graphics_spec 우선, 실패 시 visual_type 기본 렌더.
    우선순위: (1) image_path (2) ppt_shapes elements (3) apply_visual
    차트 PNG는 objective가 있으면 좌측 맥락 + 우측 이미지 2열로 가독성을 높인다.
    """
    gs = slide_spec.graphics_spec
    if not isinstance(gs, dict):
        apply_visual(slide, slide_spec)
        return

    rel = gs.get("image_path") or gs.get("asset_path")
    if rel:
        full = Path(rel)
        if not full.is_file():
            full = _project_root() / str(rel).replace("\\", "/").lstrip("/")
        if full.is_file():
            try:
                obj = (slide_spec.objective or "").strip()
                cap_w = Inches(6.05)
                gap = Inches(0.22)
                pic_left = _LEFT
                pic_w = _W
                pic_h = _IMG_H
                itop = int(_IMG_TOP)
                if len(obj) > 20:
                    cap = slide.shapes.add_textbox(_LEFT, itop, int(cap_w), int(_IMG_H))
                    tf = cap.text_frame
                    tf.word_wrap = True
                    tf.clear()
                    p1 = tf.paragraphs[0]
                    p1.text = "슬라이드 맥락 (objective)"
                    p1.font.bold = True
                    p1.font.size = Pt(11)
                    p1.font.color.rgb = theme.TITLE_RGB
                    p2 = tf.add_paragraph()
                    p2.text = obj[:1200]
                    p2.font.size = Pt(10)
                    p2.font.color.rgb = theme.BODY_RGB
                    p2.space_before = Pt(6)
                    pic_left = _LEFT + cap_w + gap
                    pic_w = _W - cap_w - gap
                slide.shapes.add_picture(str(full), int(pic_left), itop, width=int(pic_w), height=int(pic_h))
                return
            except Exception:
                pass

    if gs.get("elements") and str(gs.get("render_mode") or "") in ("ppt_shapes", "hybrid"):
        try:
            add_graphics_elements(slide, list(gs["elements"]))
            return
        except Exception:
            pass
        fb = gs.get("fallback_elements")
        if isinstance(fb, list) and fb:
            try:
                add_graphics_elements(slide, fb)
                return
            except Exception:
                pass

    apply_visual(slide, slide_spec)
