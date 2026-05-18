"""레이아웃 유틸 — python-pptx로 표·도형·텍스트를 배치한다(실패 시 Composer가 fallback 슬라이드로 대체 가능)."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from autopm.ppt import theme

# 16:9 기준 여백 — 장표가 답답하지 않게 잡는다.
_LEFT = Inches(0.55)
_RIGHT = Inches(0.55)
_TOP = Inches(0.45)
_W = Inches(12.67)  # usable width approx
# 핵심 메시지 밴드 아래에서 본문·도표 시작 — 제목/콜아웃과 겹침 방지.
CONTENT_TOP = Inches(2.02)


def add_slide_accent_strip(slide: Slide) -> None:
    """좌측 세로 악센트 — 슬라이드 첫 도형으로 넣어 다른 요소 뒤에 깔린다(발표용 톤업)."""
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.11), Inches(7.52))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.ACCENT
    bar.line.fill.background()


def add_title(slide: Slide, title: str, subtitle: str | None = None) -> None:
    """제목 상단 박스 — 남색 강조."""
    box = slide.shapes.add_textbox(_LEFT, _TOP, _W, Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme.TITLE_RGB
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = theme.SUBTITLE_RGB
        p2.space_before = Pt(6)


def add_key_message(slide: Slide, message: str) -> None:
    """슬라이드당 핵심 메시지 — 밝은 밴드 안에 넣어 한눈에 들어오게 한다."""
    if not message.strip():
        return
    band_top = Inches(1.26)
    band_h = Inches(0.62)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, _LEFT, band_top, _W, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = theme.KEY_BAND_FILL
    band.line.color.rgb = theme.ACCENT_MUTED
    tf = band.text_frame
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(8)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"핵심 메시지: {message}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = theme.KEY_MSG_RGB


def add_summary_cards(slide: Slide, cards: list[dict[str, Any]]) -> None:
    """요약 카드 — 최대 4개 가로 배치(16:9에서 발표 가독성)."""
    n = min(len(cards), 4)
    if n == 0:
        return
    gap = Inches(0.25)
    card_w = (_W - gap * (n - 1)) / n
    y = CONTENT_TOP
    h = Inches(4.45)
    for i, card in enumerate(cards[:n]):
        x = _LEFT + i * (card_w + gap)
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, card_w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = theme.CARD_FILL
        shp.line.color.rgb = theme.ACCENT
        title = str(card.get("title", f"항목 {i+1}"))
        body = str(card.get("body", ""))
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = theme.TITLE_RGB
        if body:
            p2 = tf.add_paragraph()
            p2.text = body
            p2.font.size = Pt(11)
            p2.font.color.rgb = theme.BODY_RGB
            p2.space_before = Pt(8)


def add_slide_bullets(slide: Slide, bullets: list[str], *, top: Any | None = None) -> None:
    """슬라이드 하단 bullet — 본문 도형과 함께 요점을 넣어 빈 장을 방지한다."""
    lines = [str(b).strip() for b in bullets if str(b).strip()][:5]
    if not lines:
        return
    y_top = float(Inches(6.28)) if top is None else float(top)
    tb = slide.shapes.add_textbox(_LEFT, int(y_top), _W, Inches(1.12))
    tb.text_frame.word_wrap = True
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"[:500]
        p.font.size = Pt(10)
        p.font.color.rgb = theme.BODY_RGB
        p.space_before = Pt(3) if i else Pt(0)


def add_problem_structured_cards(slide: Slide, items: list[dict[str, Any]]) -> None:
    """구조화된 문제 카드 — 2x2, 제목·설명·영향."""
    items = [x for x in items if isinstance(x, dict)][:4]
    if not items:
        add_problem_cards(slide, ["(내용 없음)"])
        return
    gap_x, gap_y = Inches(0.2), Inches(0.18)
    cols = 2
    cell_w = (_W - gap_x) / cols
    cell_h = Inches(1.92)
    y0 = float(CONTENT_TOP)
    for idx, it in enumerate(items):
        r, c = divmod(idx, cols)
        x = float(_LEFT) + c * (float(cell_w) + float(gap_x))
        y = y0 + r * (float(cell_h) + float(gap_y))
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(x),
            int(y),
            int(cell_w - Inches(0.06)),
            int(cell_h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = theme.CARD_FILL
        shp.line.color.rgb = theme.ACCENT
        tit = str(it.get("title", "이슈"))[:80]
        desc = str(it.get("description", ""))[:300]
        imp = str(it.get("impact", ""))[:100]
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = tit
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = theme.TITLE_RGB
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(9)
            p2.font.color.rgb = theme.BODY_RGB
            p2.space_before = Pt(4)
        if imp:
            p3 = tf.add_paragraph()
            p3.text = f"영향: {imp}"
            p3.font.size = Pt(9)
            p3.font.italic = True
            p3.font.color.rgb = theme.SUBTITLE_RGB
            p3.space_before = Pt(4)


def add_compact_kpis_below(slide: Slide, kpis: list[dict[str, Any]], *, y_start: Any | None = None) -> None:
    """예산 표 아래에 KPI를 작은 카드로 배치."""
    kpis = [k for k in kpis if isinstance(k, dict)][:4]
    if not kpis:
        return
    y0 = float(Inches(5.05)) if y_start is None else float(y_start)
    gap = Inches(0.12)
    n = len(kpis)
    card_w = (_W - gap * (n - 1)) / n
    h = Inches(0.88)
    for i, k in enumerate(kpis):
        x = float(_LEFT) + i * (float(card_w) + float(gap))
        body = f"{k.get('current','')} → {k.get('target','')}\n{k.get('effect','')}"[:200]
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(x),
            int(y0),
            int(card_w - Inches(0.04)),
            int(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xEE, 0xF6, 0xEE)
        shp.line.color.rgb = theme.ACCENT
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(k.get("name", "KPI"))[:60]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = theme.TITLE_RGB
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(8)
        p2.font.color.rgb = theme.BODY_RGB
        p2.space_before = Pt(2)


def add_problem_cards(slide: Slide, problems: list[str]) -> None:
    """문제 카드 — 리스트를 카드형으로 펼친다."""
    cards = [{"title": f"문제 {i+1}", "body": t} for i, t in enumerate(problems[:3])]
    if not cards:
        cards = [{"title": "문제", "body": "(내용 없음)"}]
    add_summary_cards(slide, cards)


def add_process_flow(slide: Slide, steps: list[str]) -> None:
    """간단 프로세스 — 박스 + 화살표."""
    steps = [s for s in steps if s.strip()][:5]
    if not steps:
        steps = ["단계1", "단계2"]
    n = len(steps)
    box_w = (_W - Inches(0.2) * (n - 1)) / n
    y = CONTENT_TOP + Inches(0.2)
    h = Inches(1.05)
    for i, step in enumerate(steps):
        x = _LEFT + i * (box_w + Inches(0.15))
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_w - Inches(0.12), h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step[:120]
        p.font.size = Pt(11)
        p.font.color.rgb = theme.BODY_RGB
        p.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            ax = x + box_w - Inches(0.08)
            ay = y + h / 2 - Inches(0.08)
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, ax, ay, Inches(0.22), Inches(0.16))
            arr.fill.solid()
            arr.fill.fore_color.rgb = theme.ACCENT


def add_before_after(slide: Slide, before: list[str], after: list[str]) -> None:
    """2열 Before / After 비교."""
    mid = Inches(0.3)
    half = (_W - mid) / 2
    y = CONTENT_TOP
    h = Inches(4.55)
    for label, items, off in (
        ("AS-IS / Before", before, _LEFT),
        ("TO-BE / After", after, _LEFT + half + mid),
    ):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, off, y, half, h)
        box.fill.solid()
        box.fill.fore_color.rgb = theme.CARD_FILL
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = theme.TITLE_RGB
        for it in items[:6]:
            p2 = tf.add_paragraph()
            p2.text = f"• {it}"[:200]
            p2.font.size = Pt(11)
            p2.font.color.rgb = theme.BODY_RGB
            p2.space_before = Pt(4)


def add_scope_matrix(slide: Slide, included: list[str], excluded: list[str]) -> None:
    """포함/제외 범위 2열."""
    add_before_after(slide, included or ["(포함 항목)"], excluded or ["(제외 항목)"])


def add_wbs_table(slide: Slide, rows: list[dict[str, Any]]) -> None:
    """WBS 표 — 최대 6행."""
    cols = ["단계", "작업", "기간", "담당", "산출물"]
    rcount = min(max(len(rows), 1), 6)
    ccount = len(cols)
    x, y = _LEFT, CONTENT_TOP
    wtbl, htbl = _W, Inches(0.55 * (rcount + 1))
    table = slide.shapes.add_table(rcount + 1, ccount, x, y, wtbl, htbl).table
    for j, h in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xDD, 0xE6, 0xF5)
    for i in range(rcount):
        row = rows[i] if i < len(rows) else {}
        vals = [
            str(row.get("phase", row.get("단계", f"{i+1}"))),
            str(row.get("task", row.get("작업", ""))),
            str(row.get("duration", row.get("기간", ""))),
            str(row.get("owner", row.get("담당", ""))),
            str(row.get("deliverable", row.get("산출물", ""))),
        ]
        for j, val in enumerate(vals):
            table.cell(i + 1, j).text = val[:80]


def add_budget_table(slide: Slide, rows: list[dict[str, Any]]) -> None:
    """예산 표."""
    rcount = min(max(len(rows), 1), 6)
    ccount = 3
    x, y = _LEFT, CONTENT_TOP
    wtbl, htbl = _W, Inches(0.5 * (rcount + 1))
    table = slide.shapes.add_table(rcount + 1, ccount, x, y, wtbl, htbl).table
    headers = ["항목", "예상 비용", "설명"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = RGBColor(0xDD, 0xE6, 0xF5)
    for i in range(rcount):
        row = rows[i] if i < len(rows) else {}
        table.cell(i + 1, 0).text = str(row.get("item", row.get("항목", "")))[:60]
        table.cell(i + 1, 1).text = str(row.get("cost", row.get("예상 비용", "")))[:40]
        table.cell(i + 1, 2).text = str(row.get("description", row.get("설명", "")))[:120]


def add_kpi_cards(slide: Slide, kpis: list[dict[str, Any]]) -> None:
    """KPI 숫자 카드."""
    cards = []
    for k in kpis[:4]:
        cards.append(
            {
                "title": str(k.get("name", "KPI")),
                "body": f"현재: {k.get('current', '')} → 목표: {k.get('target', '')}",
            }
        )
    add_summary_cards(slide, cards)


def add_risk_matrix(slide: Slide, risks: list[dict[str, Any]]) -> None:
    """리스크 매트릭스 표."""
    rcount = min(max(len(risks), 1), 6)
    x, y = _LEFT, CONTENT_TOP
    wtbl, htbl = _W, Inches(0.48 * (rcount + 1))
    table = slide.shapes.add_table(rcount + 1, 4, x, y, wtbl, htbl).table
    headers = ["리스크", "가능성", "영향", "대응"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = RGBColor(0xF8, 0xE1, 0xE1)
    for i in range(rcount):
        row = risks[i] if i < len(risks) else {}
        table.cell(i + 1, 0).text = str(row.get("risk", ""))[:70]
        table.cell(i + 1, 1).text = str(row.get("probability", ""))[:20]
        table.cell(i + 1, 2).text = str(row.get("impact", ""))[:20]
        table.cell(i + 1, 3).text = str(row.get("response", row.get("mitigation", "")))[:100]


def add_conclusion_box(slide: Slide, text: str) -> None:
    """결론 강조 박스."""
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, _LEFT, CONTENT_TOP + Inches(0.15), _W, Inches(3.75))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE8)
    shp.line.color.rgb = theme.ACCENT
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.strip() or "결론 및 요청사항을 정리하세요."
    p.font.size = Pt(14)
    p.font.color.rgb = theme.BODY_RGB


def add_comparison_table(slide: Slide, headers: list[str], rows: list[list[str]]) -> None:
    """옵션 대비 표 — graphics_spec 없을 때 comparison_table visual_type 폴백."""
    hdr = headers[:6] if headers else ["항목", "AS-IS", "TO-BE"]
    rcount = min(max(len(rows), 1), 7)
    ccount = len(hdr)
    x, y = _LEFT, CONTENT_TOP
    wtbl, htbl = _W, Inches(0.45 * (rcount + 1))
    table = slide.shapes.add_table(rcount + 1, ccount, x, y, wtbl, htbl).table
    for j, h in enumerate(hdr):
        table.cell(0, j).text = str(h)[:40]
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = RGBColor(0xDD, 0xE6, 0xF5)
    for i in range(rcount):
        cells = rows[i] if i < len(rows) else []
        for j in range(ccount):
            v = str(cells[j] if j < len(cells) else "")[:100]
            table.cell(i + 1, j).text = v


def add_funnel_stages(slide: Slide, stages: list[str]) -> None:
    """간단 퍼널 — TRAPEZOID 세로 스택(장표 MVP)."""
    stages = [s.strip() for s in stages if s.strip()][:6]
    if not stages:
        stages = ["인식", "검증", "적용"]
    y = float(CONTENT_TOP)
    for i, stg in enumerate(stages):
        frac = (len(stages) - i) / (len(stages) + 1)
        w = float(_W) * max(0.38, frac)
        off = (float(_W) - w) / 2
        h = Inches(0.68)
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.TRAPEZOID,
            int(_LEFT + off),
            int(y),
            int(w),
            int(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)
        shp.line.color.rgb = theme.ACCENT
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = stg[:80]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        y += float(h) + float(Inches(0.08))


def add_priority_quadrant(slide: Slide, items: list[str]) -> None:
    """2x2 우선순위 매트릭스 뼈대 + bullet — 정교한 scatter 대신 발표용 단순화."""
    gap = Inches(0.2)
    half = (_W - gap) / 2
    y = CONTENT_TOP
    h = Inches(2.25)
    labels = (
        "High Impact / Low Effort",
        "High Impact / High Effort",
        "Low Impact / Low Effort",
        "Low Impact / High Effort",
    )
    for idx, (ox, oy) in enumerate(((0, 0), (1, 0), (0, 1), (1, 1))):
        x = _LEFT + ox * (half + gap)
        yy = y + oy * (h + gap)
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, int(x), int(yy), int(half), int(h))
        box.fill.solid()
        box.fill.fore_color.rgb = theme.CARD_FILL
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = labels[idx]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = theme.TITLE_RGB
    note = "; ".join(items[:6]) if items else "(우선 과제를 배치하세요)"
    tb = slide.shapes.add_textbox(_LEFT, Inches(6.15), _W, Inches(1.0))
    tb.text_frame.text = note[:500]
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = theme.BODY_RGB


def add_org_role_map(slide: Slide, roles: list[dict[str, Any]]) -> None:
    """역할 카드 — 이름·책임을 요약 카드로 표시한다."""
    cards = []
    for r in roles[:4]:
        if isinstance(r, dict):
            cards.append(
                {"title": str(r.get("name", r.get("role", "역할"))), "body": str(r.get("responsibility", ""))[:200]}
            )
    if not cards:
        cards = [{"title": "PM", "body": "추진·조율"}, {"title": "현업", "body": "룰 정의"}]
    add_summary_cards(slide, cards)


def add_swimlane_simple(slide: Slide, lanes: list[dict[str, Any]]) -> None:
    """swimlane 복잡도 회피 — 레인별 한 줄 요약을 세로 박스로 표시(MVP)."""
    lines: list[str] = []
    for ln in lanes[:5]:
        if isinstance(ln, dict):
            steps = ", ".join(map(str, ln.get("steps", [])[:4]))
            lines.append(f"{ln.get('name', 'Lane')}: {steps}")
    if not lines:
        lines = ["현업: 입력·검증", "IT: 스크립트·배포"]
    add_before_after(slide, lines[:3], lines[3:6] if len(lines) > 3 else ["(후속 단계)"])


def add_architecture_blocks(slide: Slide, layers: list[str]) -> None:
    """아키텍처 블록 세로 스택 — PNG 실패 시 도형 폴백."""
    if not layers:
        layers = ["채널", "서비스", "데이터"]
    yy = float(CONTENT_TOP)
    for lay in layers[:5]:
        h = Inches(0.78)
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, _LEFT, int(yy), _W, int(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(lay)[:120]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        yy += float(h) + float(Inches(0.12))
