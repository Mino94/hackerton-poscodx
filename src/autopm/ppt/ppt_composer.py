"""SlideDeckSpec → project_plan.pptx — AGENTS.md 필수 산출물."""

from __future__ import annotations

import traceback
from pathlib import Path
from typing import Any

from pptx import Presentation

from autopm.ppt import layout_engine
from autopm.ppt.slide_schema import SlideDeckSpec, SlideSpec
from autopm.ppt.graphics_layout import apply_graphics_or_visual

EMU_PER_INCH = 914400


def build_fallback_slide_deck(project_title: str, subtitle: str = "Demo Mode") -> SlideDeckSpec:
    """API/파싱 실패 시에도 최소 10장을 보장하는 덱 — 외부 모듈에서 재사용한다."""
    return _sample_deck_spec(project_title, subtitle)


def create_project_plan_ppt(deck_spec: dict[str, Any], output_path: str = "outputs/project_plan.pptx") -> str:
    """
    deck_spec을 기반으로 PPTX를 만들고 경로를 반환한다.
    파싱/렌더 실패 시 최소 10장 fallback 덱으로 다시 시도한다.
    """
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    spec: SlideDeckSpec | None = None
    try:
        spec = SlideDeckSpec.model_validate(deck_spec)
    except Exception:
        spec = None
    if spec is None or not spec.slides:
        spec = _minimal_fallback_deck_from_dict(deck_spec)

    try:
        _write_pptx(spec, str(out))
        return str(out.resolve())
    except Exception:
        # python-pptx 오류 시에도 파일은 남겨 데모가 끊기지 않게 한다.
        fb = _sample_deck_spec(
            project_title=str(deck_spec.get("project_title", "AutoPM Fallback")),
            subtitle="PPTX 생성 중 오류 — fallback 적용",
        )
        try:
            _write_pptx(fb, str(out))
        except Exception:
            traceback.print_exc()
        return str(out.resolve())


def _write_pptx(spec: SlideDeckSpec, path: str) -> None:
    prs = Presentation()
    prs.slide_width = int(13.333333333333334 * EMU_PER_INCH)
    prs.slide_height = int(7.5 * EMU_PER_INCH)
    for slide_spec in spec.slides:
        layout = prs.slide_layouts[6]  # blank — 도형만 직접 배치
        slide = prs.slides.add_slide(layout)
        # 좌측 악센트는 맨 아래 레이어가 되도록 가장 먼저 추가한다.
        layout_engine.add_slide_accent_strip(slide)
        layout_engine.add_title(slide, slide_spec.title, spec.subtitle if slide_spec.slide_no == 1 else None)
        layout_engine.add_key_message(slide, slide_spec.key_message or slide_spec.objective)
        apply_graphics_or_visual(slide, slide_spec)
    prs.save(path)


def _minimal_fallback_deck_from_dict(raw: dict[str, Any]) -> SlideDeckSpec:
    title = str(raw.get("project_title") or raw.get("title") or "AutoPM")
    return _sample_deck_spec(title, "Demo / 복구 모드")


def _sample_deck_spec(project_title: str, subtitle: str) -> SlideDeckSpec:
    """AGENTS.md 필수 10장 구조를 샘플 데이터로 채운다 — API 없이도 발표 가능한 뼈대."""
    return SlideDeckSpec(
        project_title=project_title,
        subtitle=subtitle,
        slides=[
            SlideSpec(
                slide_no=1,
                title="Executive Summary",
                objective="경영진에게 과제·기간·기대효과를 한 장으로 전달한다.",
                key_message="업무 개선 과제를 4주 파일럿으로 신속히 검증한다.",
                visual_type="summary_cards",
                content={
                    "cards": [
                        {"title": "과제", "body": project_title},
                        {"title": "기간", "body": "4주 파일럿(가정)"},
                        {"title": "기대", "body": "검증 리드타임 단축·오류 조기 탐지·감사 추적성 강화"},
                    ]
                },
            ),
            SlideSpec(
                slide_no=2,
                title="현재 문제점",
                objective="왜 지금 개선이 필요한지 공감대를 형성한다.",
                key_message="수작업·기준 편차로 품질 리스크가 있다.",
                visual_type="problem_cards",
                content={"problems": ["검증 리드타임 장기화", "담당자별 기준 불일치", "누락 가능성"]},
            ),
            SlideSpec(
                slide_no=3,
                title="AS-IS 프로세스",
                objective="현재 데이터 흐름과 병목을 단계적으로 보여 준다.",
                key_message="ERP→엑셀→수동 검증 흐름이다.",
                visual_type="process_flow",
                content={"steps": ["데이터 추출", "엑셀 취합", "수동 검증", "이슈 조치"]},
            ),
            SlideSpec(
                slide_no=4,
                title="TO-BE 프로세스",
                objective="개선 후 자동화·표준화 포인트를 대비 구조로 설명한다.",
                key_message="규칙 표준화와 자동 검증으로 전환한다.",
                visual_type="before_after",
                content={
                    "before": ["수동 검증", "분산 기준"],
                    "after": ["자동 룰 검증", "중앙 규칙/로그"],
                },
            ),
            SlideSpec(
                slide_no=5,
                title="개발 범위",
                objective="MVP 내·외를 명확히 하여 승인 범위 리스크를 줄인다.",
                key_message="MVP는 자동 검증·리포트에 집중한다.",
                visual_type="scope_matrix",
                content={
                    "included": ["검증 룰 MVP", "리포트/로그"],
                    "excluded": ["ERP 커스터", "대규모 인프라"],
                },
            ),
            SlideSpec(
                slide_no=6,
                title="WBS / 추진 일정",
                objective="4주 내 실행 가능한 마일스톤과 산출물을 제시한다.",
                key_message="킥오프→규칙→구현→파일럿 순으로 진행한다.",
                visual_type="wbs_table",
                content={
                    "rows": [
                        {
                            "phase": "1",
                            "task": "킥오프",
                            "duration": "3일",
                            "owner": "PM",
                            "deliverable": "워크숍 메모",
                        },
                        {
                            "phase": "2",
                            "task": "규칙 정리",
                            "duration": "1주",
                            "owner": "현업",
                            "deliverable": "룰 명세",
                        },
                        {
                            "phase": "3",
                            "task": "구현",
                            "duration": "2주",
                            "owner": "IT",
                            "deliverable": "스크립트",
                        },
                        {
                            "phase": "4",
                            "task": "파일럿",
                            "duration": "1주",
                            "owner": "전체",
                            "deliverable": "결과 리포트",
                        },
                    ]
                },
            ),
            SlideSpec(
                slide_no=7,
                title="예산 및 ROI",
                objective="비용 항목과 정성·정량 효과를 (가정) 명시해 투자 대비 납득을 돕는다.",
                key_message="비용·절감은 모두 가정 기반으로 표기한다.",
                visual_type="budget_table",
                content={
                    "rows": [
                        {"item": "분석/구현", "cost": "300만 원(가정)", "description": "인력·도구"},
                        {"item": "절감(가정)", "cost": "월 40h", "description": "검증 자동화·운영 시간"},
                    ]
                },
            ),
            SlideSpec(
                slide_no=8,
                title="리스크 및 대응",
                objective="상위 리스크와 완화책을 한 표로 정리한다.",
                key_message="데이터 품질·규칙 불완전을 상위 리스크로 관리한다.",
                visual_type="risk_matrix",
                content={
                    "risks": [
                        {
                            "risk": "데이터 품질",
                            "probability": "중",
                            "impact": "중",
                            "response": "샘플 검증",
                        },
                        {
                            "risk": "규칙 누락",
                            "probability": "중",
                            "impact": "고",
                            "response": "파일럿 2회",
                        },
                    ]
                },
            ),
            SlideSpec(
                slide_no=9,
                title="기대효과",
                objective="KPI로 파일럿 성과를 측정·확대 적용할지 판단한다.",
                key_message="시간·품질 KPI로 효과를 추적한다.",
                visual_type="kpi_cards",
                content={
                    "kpis": [
                        {"name": "검증 리드타임", "current": "기준선", "target": "-30%(가정)"},
                        {"name": "누락 건수", "current": "미측정", "target": "0건 지향"},
                    ]
                },
            ),
            SlideSpec(
                slide_no=10,
                title="결론 / 요청사항",
                objective="다음 액션·승인 요청을 명확히 남긴다.",
                key_message="RACI·샘플 데이터·보안 합의를 요청한다.",
                visual_type="conclusion_box",
                content={
                    "text": "승인을 위해 파일럿 범위·데이터 접근·일정을 확정해 주세요. RACI와 검증 스코프에 합의가 필요합니다."
                },
            ),
        ],
    )
